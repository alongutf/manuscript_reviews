"""
Build a 12-minute physics-audience talk deck (PPTX) for the dysregulated-persistence
project. Theme: collapsing a 10^4-dimensional regulatory network onto a single
coupling-strength parameter rho (SK spin-glass analogy).

Embeds figures already in the repo (PNG previews). Run from anywhere:
    python scripts/build_presentation.py
Writes: presentations/dysregulated_persistence_physics_talk.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ── paths ────────────────────────────────────────────────────────────────────
HERE = os.path.dirname(os.path.abspath(__file__))
REPO = os.path.dirname(HERE)
OUT_DIR = os.path.join(REPO, "presentations")
os.makedirs(OUT_DIR, exist_ok=True)
OUT = os.path.join(OUT_DIR, "dysregulated_persistence_physics_talk.pptx")

def P(*parts):
    """Join path parts onto the repo root."""
    return os.path.join(REPO, *parts)

# figure PNGs already produced by the figureN.py / simulation scripts; keys are the
# slide "image_key" values used by content_slide() below
FIG = {
    "intro":     P("scripts", "figures", "figure1", "experiment illustration.png"),
    "theory":    P("scripts", "figures", "figure2_preview.png"),
    "simulator": P("scripts", "supplementary_figures", "figure_s2_preview.png"),
    "calib":     P("results", "simulation_results", "figures", "rho_sweep_calibration.png"),
    "climax":    P("scripts", "figures", "figure3_preview.png"),
    "robust":    P("results", "simulation_results", "figures",
                   "subpopulation_mixing_rho09_50_50_20260615_111850.png"),
    "biology":   P("scripts", "figures", "figure4_preview.png"),
}

# ── colors / theme ───────────────────────────────────────────────────────────
INK   = RGBColor(0x1A, 0x1A, 0x2E)   # near-black navy
ACCENT= RGBColor(0x4C, 0x72, 0xB0)   # steelblue  (regulated / high rho)
WARM  = RGBColor(0xE0, 0x7B, 0x54)   # coral      (dysregulated / low rho)
MUTE  = RGBColor(0x55, 0x55, 0x66)
BG    = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ── helpers ──────────────────────────────────────────────────────────────────
def add_slide():
    """New blank slide (BLANK layout, no placeholders) with a solid white background."""
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s

def textbox(s, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, bold, color, space_after_pt)."""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, size, bold, color, sa) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sa)
        r = p.add_run()
        r.text = txt
        f = r.font
        f.size = Pt(size); f.bold = bold; f.color.rgb = color
        f.name = "Calibri"
    return tb

def bar(s, color=ACCENT, h=Inches(0.12)):
    """Full-width flat color bar across the top of the slide (autoshape type 1 =
    MSO_SHAPE.RECTANGLE); used as the accent stripe under the title band and as the
    top/bottom frame bars on the title slide."""
    sp = s.shapes.add_shape(1, 0, 0, SW, h)  # rectangle
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    return sp

def title_band(s, title, kicker=None):
    """Accent bar + optional small caps "kicker" label + slide title, top-left."""
    bar(s, ACCENT)
    lines = []
    if kicker:
        lines.append((kicker.upper(), 13, True, ACCENT, 2))
    lines.append((title, 30, True, INK, 0))
    textbox(s, Inches(0.55), Inches(0.28), Inches(12.2), Inches(1.1), lines)

def fit_image(s, path, box_l, box_t, box_w, box_h, align="center"):
    """Place image scaled to fit inside box, preserving aspect ratio."""
    if not os.path.exists(path):
        textbox(s, box_l, box_t, box_w, box_h,
                [("[missing figure: %s]" % os.path.basename(path), 14, False, WARM, 0)],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return None
    iw, ih = Image.open(path).size
    ar = iw / ih
    bw, bh = box_w, box_h
    box_ar = bw / bh
    # letterbox to the tighter dimension: if the image is relatively wider than the
    # box, width is the binding constraint (and height is derived), else height is
    if ar > box_ar:
        w = bw; h = int(bw / ar)
    else:
        h = bh; w = int(bh * ar)
    # center the scaled image inside the box on both axes
    l = box_l + (bw - w) // 2
    t = box_t + (bh - h) // 2
    return s.shapes.add_picture(path, l, t, width=Emu(int(w)), height=Emu(int(h)))

def caption(s, l, t, w, txt, color=MUTE):
    """Small centered caption line, e.g. below a figure."""
    textbox(s, l, t, w, Inches(0.5),
            [(txt, 13, False, color, 0)], align=PP_ALIGN.CENTER)

def content_slide(title, kicker, image_key, bullets, cap=None, img_side="right"):
    """Standard split slide: figure on one side, bullets on the other."""
    s = add_slide()
    title_band(s, title, kicker)
    img_w = Inches(7.1); img_h = Inches(5.2)
    txt_w = Inches(4.5)
    # NOTE: despite the parameter name, img_side="right" places the image on the LEFT
    # (img_l=0.55) with text on the right, and img_side="left" places the image on the
    # right-hand portion of the slide (img_l=5.75) with text on the left -- the two
    # branches are swapped relative to what "img_side" suggests. See log file for
    # this script (logs/scripts/build_presentation.txt) for details; left as-is here
    # since this is a comments-only documentation pass.
    if img_side == "right":
        img_l = Inches(0.55); txt_l = Inches(8.0)
    else:
        img_l = Inches(5.75); txt_l = Inches(0.55)
    top = Inches(1.6)
    fit_image(s, FIG[image_key], img_l, top, img_w, img_h)
    if cap:
        caption(s, img_l, top + img_h + Inches(0.02), img_w, cap)
    lines = []
    for b, sub in bullets:
        lines.append(("•  " + b, 19, True, INK, 6))
        if sub:
            lines.append(("    " + sub, 15, False, MUTE, 12))
    textbox(s, txt_l, top, txt_w, img_h, lines, anchor=MSO_ANCHOR.TOP)
    return s

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
bar(s, ACCENT, Inches(0.18))
bar2 = s.shapes.add_shape(1, 0, SH - Inches(0.18), SW, Inches(0.18))
bar2.fill.solid(); bar2.fill.fore_color.rgb = WARM; bar2.line.fill.background()
textbox(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(2.8), [
    ("From 10⁴ genes to one number", 44, True, INK, 6),
    ("A coupling-strength order parameter ρ for cellular dysregulation", 24, False, ACCENT, 18),
    ("Genome-wide dysregulation in antibiotic tolerance and persistence", 17, False, MUTE, 0),
])
textbox(s, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.8), [
    ("Random-matrix theory  ·  Gaussian-copula simulations  ·  bacterial scRNA-seq", 15, False, MUTE, 0),
])

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — The physics question / SK analogy
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
title_band(s, "The cell as a strongly-interacting many-body system", "the question")
textbox(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(5.4), [
    ("A transcriptome is ~10³–10⁴ genes coupled by a regulatory network.", 22, True, INK, 14),
    ("The spin-glass question:", 20, True, ACCENT, 6),
    ("Sherrington–Kirkpatrick — N spins, all-to-all random couplings Jᵢⱼ. The macroscopic "
     "state is set by the width J of the coupling distribution, not the ∼N² individual Jᵢⱼ.",
     18, False, INK, 14),
    ("Our claim:", 20, True, WARM, 6),
    ("the genome-wide regulatory state collapses onto a single coupling strength ρ "
     "(shared-variance fraction) — analogous to J. The whole talk: define ρ, recover it from "
     "simulations, measure it in data.", 18, False, INK, 14),
    ("Biology, stated once: stress (antibiotics) acts like raising temperature — it melts the "
     "ordered network. Tolerant/persister cells = low ρ.", 15, False, MUTE, 0),
])

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Setup: data as a random matrix
# ─────────────────────────────────────────────────────────────────────────────
content_slide(
    "The data is a large random matrix", "setup", "intro",
    [("Count matrix X : n cells × p genes", "p ∼ 10³–10⁴, n ∼ 10³  →  ratio γ = p/n = O(1)"),
     ("Object of interest: the gene–gene correlation matrix", "C = (1/n) XᵀX  after row-normalize + z-score"),
     ("This is the RMT regime", "sample correlations are dominated by noise — we need the right null"),
     ("Plan", "(1) noise null, (2) order parameter, (3) calibrate with simulations, (4) measure data")],
    cap="Single-cell RNA-seq of bacteria under acute stress (experimental setup).",
    img_side="left",
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Theory: MP null + Generalized MP + order parameter  (Figure 2)
# ─────────────────────────────────────────────────────────────────────────────
content_slide(
    "Null model, signal, and an order parameter", "theory",
    "theory",
    [("Marchenko–Pastur null", "independent genes → eigenvalues fill the MP sea λ± = (1±√γ)² (the ρ=0 reference)"),
     ("Realize it by scrambling", "permute each gene across cells — kills correlations, keeps marginals; matches MP"),
     ("Correlations push eigenvalues out", "outliers above λ_max^scr = BBP-type transition"),
     ("Generalized MP (GMP)", "a one-parameter ρ deformation of the bulk — fit ρ to the whole spectrum"),
     ("Order parameter:", "GMP-Cor = Σ max(λᵢ − λ_max^scr, 0) = spectral weight expelled above the noise edge")],
    cap="Fig 2: random-matrix MP, generalized-MP at ρ=0/0.7/0.8/0.9, eigenvalue density & data CCDFs.",
    img_side="right",
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — The generative model behind rho  (text + small equation feel)
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
title_band(s, "What ρ means microscopically", "the model")
textbox(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(5.4), [
    ("Factor model for the ground-truth correlation matrix:", 21, True, INK, 12),
    ("A = √(1−ρ)·I   +   cluster factors   +   one global hub factor      ;     C = AAᵀ, normalized to unit diagonal",
     18, True, ACCENT, 14),
    ("ρ = shared-variance fraction — it sets the mean off-diagonal magnitude (the coupling width).",
     18, False, INK, 14),
    ("Two-tier structure mirrors the biology at low resolution:", 18, True, WARM, 6),
    ("clusters ≈ operons / regulons   ·   global hub ≈ master regulator.", 17, False, INK, 14),
    ("ρ → 1 : tightly coupled, ordered network   ·   ρ → 0 : decoupled, dysregulated.",
     18, True, INK, 0),
])

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — The simulator  (Figure S2)
# ─────────────────────────────────────────────────────────────────────────────
content_slide(
    "A realistic instrument with known ρ", "simulations",
    "simulator",
    [("Goal", "test whether GMP-Cor recovers the microscopic coupling under realistic noise"),
     ("Stage 1: structure", "build C(ρ) from the cluster+hub factor model"),
     ("Stage 2: counts via Gaussian copula", "latent z~N(0,C) → Φ → NB counts (inv-Gamma means)"),
     ("Realistic nuisances", "library-size variation + expression-dependent dropout"),
     ("Why it works", "two-stage design decouples ground-truth ρ from count-level noise")],
    cap="Fig S2: factor-model network → correlation matrix → copula count pipeline → simulated output.",
    img_side="right",
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Calibration curve (key result #1)
# ─────────────────────────────────────────────────────────────────────────────
content_slide(
    "Key result: GMP-Cor faithfully recovers ρ", "calibration",
    "calib",
    [("Sweep ρ, hold structure fixed", "only count-sampling noise varies across repeats"),
     ("GMP-Cor(ρ) is monotonic", "and robust — small error bars over replicates"),
     ("The order parameter is invertible", "one number reports the microscopic coupling strength"),
     ("This is the ‘it works’ slide", "spectral readout ↔ hidden ρ, by construction-blind calibration")],
    cap="ρ-sweep calibration: GMP-Cor increases monotonically with the coupling strength ρ.",
    img_side="left",
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Robustness / non-trivial test (subpopulation mixing)
# ─────────────────────────────────────────────────────────────────────────────
content_slide(
    "It measures coupling, not heterogeneity", "robustness",
    "robust",
    [("The obvious objection", "‘you’re just detecting clusters / batch structure’"),
     ("Stress test", "mix two subpopulations with distinct hub topologies but the same ρ"),
     ("Result", "GMP-Cor separates a mix of two regulated pops (high ρ) from two dysregulated pops (low ρ)"),
     ("Conclusion", "the metric reports coupling strength, not mere population mixing")],
    cap="50/50 subpopulation-mixing control: GMP-Cor tracks ρ, not the number of populations.",
    img_side="right",
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Real data: simulation <-> data bridge (Figure 3, the climax)
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
title_band(s, "Real cells fall on the calibration curve", "data")
fit_image(s, FIG["climax"], Inches(0.55), Inches(1.55), Inches(7.5), Inches(5.5))
caption(s, Inches(0.55), Inches(7.05), Inches(7.5),
        "Fig 3: experimental Reg- vs Dis-Arrest spectra, matched simulations (ρ=0.9 / 0.5), "
        "calibration curve, group comparison.")
textbox(s, Inches(8.3), Inches(1.7), Inches(4.6), Inches(5.4), [
    ("Empirical spectra → fitted ρ", 19, True, INK, 8),
    ("High-ρ vs low-ρ samples have visibly different eigenvalue spectra; each collapses to one ρ.",
     15, False, MUTE, 14),
    ("Simulations match data", 19, True, ACCENT, 8),
    ("ρ=0.9 sim ↔ regulated cells   ·   ρ=0.5 sim ↔ dysregulated cells.", 15, False, MUTE, 14),
    ("Groups separate on ρ", 19, True, WARM, 8),
    ("Regulated vs Dis-Arrest populations land at distinct, significant GMP-Cor — right on the "
     "calibration curve.", 15, False, MUTE, 0),
])

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Biology at low resolution (Figure 4) — one slide only
# ─────────────────────────────────────────────────────────────────────────────
content_slide(
    "The biology, at low resolution", "interpretation",
    "biology",
    [("Low ρ = persistence/tolerance", "globally dysregulated, decoupled regulatory network"),
     ("Coordinated programs break down", "e.g. coherent gene modules lose their correlation structure"),
     ("Phenotype link", "dysregulation tracks antibiotic tolerance / lag-time survival"),
     ("(Mentioned once — back to the physics)", "")],
    cap="Fig 4: downstream functional read-out of the dysregulated state.",
    img_side="left",
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Close: back to the spin glass
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
title_band(s, "One number for a 10⁴-dimensional network", "takeaway")
textbox(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(5.4), [
    ("What we bought:", 22, True, INK, 10),
    ("(1) a clean RMT null (Marchenko–Pastur via scrambling)", 18, False, INK, 6),
    ("(2) a spectral order parameter, GMP-Cor, with a microscopic meaning (ρ)", 18, False, INK, 6),
    ("(3) a calibrated, robust copula simulator linking the two", 18, False, INK, 14),
    ("The SK payoff:", 22, True, WARM, 10),
    ("just as J controls the spin glass without knowing any individual Jᵢⱼ, ρ summarizes the "
     "cell’s global regulatory coupling — and stress acts like temperature, melting the ordered "
     "network.", 19, True, INK, 14),
    ("Open question: is the loss of coordination a sharp transition in ρ?  —  thank you / questions.",
     17, False, ACCENT, 0),
])

prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
# report which figure inputs were actually found on disk; fit_image() already falls
# back to a "[missing figure: ...]" placeholder textbox at build time, so a MISSING
# line here just flags the slide(s) that shipped with a placeholder instead of an image
for k, v in FIG.items():
    print(("  OK  " if os.path.exists(v) else "  MISSING  "), k, "->", os.path.relpath(v, REPO))
